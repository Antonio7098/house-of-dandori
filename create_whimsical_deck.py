from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BRAND_COLORS = {
    "forest": RGBColor(65, 153, 115),
    "forest_dark": RGBColor(26, 61, 46),
    "terracotta": RGBColor(232, 138, 100),
    "honey": RGBColor(232, 184, 74),
    "sage": RGBColor(165, 181, 152),
    "cream": RGBColor(250, 248, 244),
}

TITLE_FONT = "Fraunces"
BODY_FONT = "Crimson Pro"
OUTPUT_PATH = "School_of_Dandori_Context_Playbook.pptx"


def add_background(slide, width, height, color, transparency=0.0):
    shapes = slide.shapes
    background = shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left=0,
        top=0,
        width=width,
        height=height,
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.fill.transparency = transparency
    background.line.fill.background()
    background.z_order = 1
    return background


def add_accent(slide, left, top, width, height, color, transparency=0.2):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.z_order = 2
    return shape


def add_title(slide, text, subtitle=None, center=False, color=RGBColor(13, 31, 23)):
    title_shape = slide.shapes.title or slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.5), Inches(1.5))
    title_tf = title_shape.text_frame
    title_tf.clear()
    paragraph = title_tf.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = TITLE_FONT
    paragraph.font.size = Pt(48)
    paragraph.font.color.rgb = color
    if center:
        paragraph.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(8.5), Inches(1.0))
        sub_tf = sub_shape.text_frame
        sub_tf.clear()
        sub_paragraph = sub_tf.paragraphs[0]
        sub_paragraph.text = subtitle
        sub_paragraph.font.name = BODY_FONT
        sub_paragraph.font.size = Pt(20)
        sub_paragraph.font.color.rgb = RGBColor(80, 84, 76)
        if center:
            sub_paragraph.alignment = PP_ALIGN.CENTER


def add_bullets(slide, bullet_points, left=Inches(0.9), top=Inches(2.3), width=Inches(8.2), height=Inches(4.5), font_color=RGBColor(26, 61, 46)):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    for idx, point in enumerate(bullet_points):
        p = text_frame.add_paragraph() if idx else text_frame.paragraphs[0]
        p.text = point
        p.font.name = BODY_FONT
        p.font.size = Pt(22)
        p.level = 0
        p.font.color.rgb = font_color


slides = [
    {
        "layout": 0,
        "title": "School of Dandori",
        "subtitle": "Preserving Whimsy with Contextually Grounded AI",
        "bg": BRAND_COLORS["cream"],
        "accents": [
            (Inches(-1), Inches(-1), Inches(4), Inches(4), BRAND_COLORS["forest"], 0.25),
            (Inches(7), Inches(-0.5), Inches(4), Inches(4), BRAND_COLORS["terracotta"], 0.2),
        ],
    },
    {
        "layout": 5,
        "title": "Agenda",
        "bullets": [
            "Why whimsy matters",
            "Visual + messaging system",
            "Vibe search & styling",
            "Context engineering stack",
            "Chatbot toolbelt & guardrails",
        ],
        "bg": BRAND_COLORS["cream"],
        "accents": [
            (Inches(6), Inches(0.3), Inches(3.5), Inches(3.5), BRAND_COLORS["honey"], 0.25),
        ],
    },
    {
        "layout": 5,
        "title": "Preserving the Whimsical Vibe",
        "bullets": [
            "Visual cues: organic gradients, botanical textures, warm serif typography",
            "Copywriting: invitations to play, ‘describe your mood’ prompts, gentle humor",
            "Micro-interactions: floating sparkles, elastic buttons, calming motion curve",
            "Dedicated /visual-elements route for full-screen capture of all UI atoms",
        ],
        "bg": BRAND_COLORS["sage"],
        "accents": [
            (Inches(-0.5), Inches(3.5), Inches(3), Inches(3), BRAND_COLORS["forest"], 0.35),
        ],
    },
    {
        "layout": 5,
        "title": "Vibe Search Experience",
        "bullets": [
            "Guided hero search plus ‘Vibe Search’ CTA that opens the chat panel",
            "Semantic graphSearch endpoint fuels intent-based matching across filters",
            "Fallback to SQL search guarantees graceful degradation when embeddings fail",
            "Chat input hints reinforce playful tone: ‘Whisper a detail or two…’",
        ],
        "bg": BRAND_COLORS["cream"],
        "accents": [
            (Inches(7), Inches(3.8), Inches(4), Inches(4), BRAND_COLORS["terracotta"], 0.25),
        ],
    },
    {
        "layout": 5,
        "title": "Styling System",
        "bullets": [
            "Color tokens mirror theme.css forest/terracotta/honey palette",
            "Fraunces for display, Crimson Pro for body; fluid type scale keeps slides elegant",
            "Cards, badges, search box rendered via production components for faithful screenshots",
            "Animations (leaf float, sparkles) curated to stay calming yet magical",
        ],
        "bg": BRAND_COLORS["sage"],
        "accents": [
            (Inches(0.5), Inches(0.2), Inches(3.5), Inches(3.5), BRAND_COLORS["honey"], 0.3),
        ],
    },
    {
        "layout": 5,
        "title": "Messaging Pillars",
        "bullets": [
            "Purpose-driven: Reclaim attention, nurture adult play, celebrate instructors",
            "Warm authority: references Ada, Tessa, Arthur—the caretakers of the school",
            "Tactile verbs: ‘craft’, ‘conjure’, ‘tend’, grounding AI copy in human action",
            "Evidence-first promises: every magical suggestion is backed by real course data",
        ],
        "bg": BRAND_COLORS["cream"],
        "accents": [
            (Inches(-0.4), Inches(0.5), Inches(3), Inches(3), BRAND_COLORS["forest"], 0.3),
        ],
    },
    {
        "layout": 5,
        "title": "Context Engineering North Star",
        "bullets": [
            "Guarantee verifiable answers: SQL → Vector RAG → GraphRAG escalation",
            "Deterministic metadata hygiene via sanitize_metadata + chunk builders",
            "Configurable providers (Chroma, Qdrant, Vertex, Neo4j) swapped via env vars",
            "Playbook inspired by docs/CONTEXT_ENGINEERING.md for onboarding clarity",
        ],
        "bg": BRAND_COLORS["sage"],
        "accents": [
            (Inches(6.5), Inches(0.1), Inches(3.5), Inches(3.5), BRAND_COLORS["forest_dark"], 0.25),
        ],
    },
    {
        "layout": 5,
        "title": "Layered Retrieval Stack",
        "bullets": [
            "SQL search_courses: fast triage + sorting via allow-listed filters",
            "Vector RAG: CourseChunkBuilder(mode='simple') with normalized _shape_results()",
            "GraphRAG: enriched triples + narrative chunks plus optional Neo4j neighbors",
            "Examples JSON prove fidelity ladder: rag_moss → graphrag_moss → neo4j_moss",
        ],
        "bg": BRAND_COLORS["cream"],
        "accents": [
            (Inches(7.2), Inches(0.4), Inches(3.2), Inches(3.2), BRAND_COLORS["honey"], 0.25),
        ],
    },
    {
        "layout": 5,
        "title": "Chatbot Toolbelt",
        "bullets": [
            "Tool-enforced retrieval: search_courses, semantic_search, graph_neighbors",
            "_initial_context seeds SQL titles + RAG summaries + optional graph neighbors",
            "Artifacts + display(id) markers hydrate UI capsules with evidence snippets",
            "Graceful degradation messaging keeps trust high when API keys or providers fail",
        ],
        "bg": BRAND_COLORS["sage"],
        "accents": [
            (Inches(-0.2), Inches(3.6), Inches(3.5), Inches(3.5), BRAND_COLORS["terracotta"], 0.3),
        ],
    },
    {
        "layout": 5,
        "title": "Next Steps",
        "bullets": [
            "Capture UI slices via /visual-elements for deck imagery",
            "Run scripts/reindex_services.py --mode both after ingestion tweaks",
            "Enable GRAPH_RAG_USE_NEO4J in staging to demo neighbor bullets live",
            "Continue evolving messaging samples and vibe prompts with instructor feedback",
        ],
        "bg": BRAND_COLORS["cream"],
        "accents": [
            (Inches(5.8), Inches(3.9), Inches(4), Inches(4), BRAND_COLORS["forest"], 0.25),
        ],
    },
]


prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
SLIDE_WIDTH = prs.slide_width
SLIDE_HEIGHT = prs.slide_height

for slide_spec in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[slide_spec.get("layout", 5)])
    background = add_background(
        slide,
        width=SLIDE_WIDTH,
        height=SLIDE_HEIGHT,
        color=slide_spec.get("bg", BRAND_COLORS["cream"]),
    )
    background.z_order = 1

    for accent in slide_spec.get("accents", []):
        add_accent(slide, *accent)

    add_title(slide, slide_spec["title"], slide_spec.get("subtitle"), center=(slide_spec["layout"] == 0))

    bullets = slide_spec.get("bullets")
    if bullets:
        add_bullets(slide, bullets)

prs.save(OUTPUT_PATH)
print(f"Deck created at {OUTPUT_PATH}")
